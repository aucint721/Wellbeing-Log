#!/usr/bin/env python3
"""
AI Presentation Generator
Creates beautiful presentations from text using local Ollama models or Claude API.
Polished designer layouts: concise bullets, tight spacing, sequential appear animations.
"""

import json
import os
import requests
import yaml
from typing import Dict, List, Optional

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# PowerPoint DrawingML / PresentationML namespaces
_NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}
_P = "{%s}" % _NSMAP["p"]


class PresentationGenerator:
    """Generate presentations using Ollama models or Claude API with polished designer layouts."""

    MAX_BULLETS = 4
    MAX_BULLET_CHARS = 80
    MAX_BULLET_WORDS = 12

    def __init__(self, config_path: str = "config.yaml"):
        with open(config_path, "r") as f:
            self.config = yaml.safe_load(f)

        self.ollama_url = self.config["ollama"]["base_url"]
        self.models = self.config["models"]
        self.themes = self.config["themes"]
        self.claude_config = self.config.get("claude", {})

    def list_models(self) -> Dict:
        """List available models with details."""
        return self.models

    def list_themes(self) -> Dict:
        """List available themes."""
        return self.themes

    def _ollama_installed_names(self) -> Optional[List[str]]:
        """Return installed Ollama model names, or None if Ollama is unreachable."""
        try:
            response = requests.get(f"{self.ollama_url}/api/tags", timeout=3)
            response.raise_for_status()
            models = response.json().get("models") or []
            names: List[str] = []
            for entry in models:
                name = (entry.get("name") or entry.get("model") or "").strip()
                if name:
                    names.append(name)
            return names
        except (requests.RequestException, ValueError, TypeError):
            return None

    @staticmethod
    def _ollama_model_pulled(model_id: str, installed: List[str]) -> bool:
        """True if model_id matches an installed Ollama tag (with/without :tag)."""
        wanted = model_id.strip().lower()
        wanted_base = wanted.split(":", 1)[0]
        for name in installed:
            n = name.strip().lower()
            if n == wanted or n.split(":", 1)[0] == wanted_base:
                return True
        return False

    def list_models_with_status(self) -> Dict:
        """
        Return models plus live availability for the Web UI.

        Claude is available when ANTHROPIC_API_KEY (or config api_key) is set.
        Ollama models need a reachable Ollama server and a pulled model image.
        """
        installed = self._ollama_installed_names()
        ollama_up = installed is not None
        claude_key = bool(os.getenv("ANTHROPIC_API_KEY") or self.claude_config.get("api_key"))

        enriched: Dict = {}
        for key, model in self.models.items():
            info = dict(model)
            provider = model.get("provider", "ollama")
            if provider == "claude":
                info["available"] = claude_key
                info["status"] = (
                    "Ready"
                    if claude_key
                    else "Set ANTHROPIC_API_KEY to use Claude"
                )
            else:
                model_id = model.get("model_id", key)
                if not ollama_up:
                    info["available"] = False
                    info["status"] = (
                        f"Ollama is not running at {self.ollama_url}. "
                        "Start it (open the Ollama app, or run: ollama serve)"
                    )
                elif self._ollama_model_pulled(model_id, installed or []):
                    info["available"] = True
                    info["status"] = "Ready (local Ollama)"
                else:
                    info["available"] = False
                    info["status"] = f"Not installed. Run: ollama pull {model_id}"
            enriched[key] = info

        return {
            "models": enriched,
            "ollama": {
                "up": ollama_up,
                "base_url": self.ollama_url,
                "installed": installed or [],
            },
            "claude": {"api_key_set": claude_key},
        }

    def _build_prompt(self, content: str, num_slides: int) -> str:
        """Build the polished concise content-generation prompt."""
        return f"""You are an expert presentation designer creating a captivating presentation from the following content.

Content:
{content}

CRITICAL CONTENT RULES:
- Each content slide must have EXACTLY 3-4 bullet points (NO MORE)
- Each bullet must be SHORT — maximum 8-10 words, one line on screen
- Make every word count — punchy, memorable, active verbs
- Include specific numbers or examples when possible

BAD (too long):
"Machine learning is a subset of artificial intelligence that enables computers to learn from data without being explicitly programmed"

GOOD (concise):
"Computers learn from data without explicit programming"
"Performance improves automatically with more examples"
"Powers Netflix recommendations and spam filters"

TARGET AUDIENCE: Year 1 students (or general audience)
- Simple, clear language
- Visual and memorable
- Fun and engaging, not dry

Create exactly {num_slides} slides with this structure:
1. Title slide (short catchy title + brief subtitle)
2. Content slides (exactly 3-4 short bullets)
3. Section breaks (1-2 for major transitions)
4. Conclusion slide

Return ONLY valid JSON in this exact format:
{{
  "title": "Short, Catchy Title (4-6 words max)",
  "slides": [
    {{
      "type": "title",
      "title": "Short Title (4-6 words)",
      "subtitle": "Brief subtitle (8-10 words max)"
    }},
    {{
      "type": "content",
      "title": "Clear Slide Title (3-5 words)",
      "bullets": [
        "First short point (8-10 words max)",
        "Second concise insight (8-10 words max)",
        "Third memorable fact (8-10 words max)"
      ]
    }},
    {{
      "type": "section",
      "title": "Section Name (2-4 words)"
    }},
    {{
      "type": "conclusion",
      "title": "Thank You"
    }}
  ]
}}

IMPORTANT:
- Keep ALL text SHORT
- 3-4 bullets per content slide (NO MORE)
- Each bullet: ONE LINE only (8-10 words)
- Be specific and memorable
- Base the presentation on the provided content

Generate the presentation outline now:"""

    def _normalize_outline(self, outline: Dict) -> Dict:
        """Enforce concise bullets so slides never overflow."""
        for slide in outline.get("slides", []):
            if slide.get("type") != "content":
                continue
            bullets = slide.get("bullets") or []
            trimmed = []
            for bullet in bullets[: self.MAX_BULLETS]:
                text = " ".join(str(bullet).split())
                words = text.split()
                if len(words) > self.MAX_BULLET_WORDS:
                    text = " ".join(words[: self.MAX_BULLET_WORDS])
                if len(text) > self.MAX_BULLET_CHARS:
                    text = text[: self.MAX_BULLET_CHARS].rstrip() + "…"
                if text:
                    trimmed.append(text)
            slide["bullets"] = trimmed
        return outline

    def _extract_json_outline(self, response_text: str) -> Dict:
        """Extract and parse JSON outline from model response text."""
        response_text = response_text.strip()
        json_start = response_text.find("{")
        json_end = response_text.rfind("}") + 1

        if json_start < 0 or json_end <= json_start:
            raise ValueError("No valid JSON found in response")

        outline = json.loads(response_text[json_start:json_end])
        if "slides" not in outline:
            raise ValueError("Outline JSON missing 'slides' key")
        return self._normalize_outline(outline)

    def _generate_outline_ollama(self, content: str, model_key: str, num_slides: int) -> Dict:
        """Generate outline via local Ollama."""
        model_id = self.models[model_key]["model_id"]
        model_name = self.models[model_key]["name"]
        prompt = self._build_prompt(content, num_slides)

        print(f"Generating outline with {model_name} (Ollama)...")

        try:
            response = requests.post(
                f"{self.ollama_url}/api/generate",
                json={
                    "model": model_id,
                    "prompt": prompt,
                    "stream": False,
                    "options": {
                        "temperature": 0.7,
                        "top_p": 0.9,
                    },
                },
                timeout=300,
            )
        except requests.exceptions.ConnectionError as e:
            raise ConnectionError(
                f"{model_name} needs Ollama running locally. "
                f"Could not connect to {self.ollama_url}. "
                "On Mac: open the Ollama app (or run `ollama serve`), then "
                f"`ollama pull {model_id}`, and try again. "
                "Or switch the model to Claude (uses ANTHROPIC_API_KEY)."
            ) from e
        except requests.exceptions.Timeout as e:
            raise TimeoutError(
                f"Ollama timed out while generating with {model_name} ({model_id}). "
                "Large models can be slow on first load — wait for the model to finish "
                "loading in Ollama, or try Dolphin 8B / Claude."
            ) from e

        if response.status_code != 200:
            body = (response.text or "")[:300]
            lower = body.lower()
            if response.status_code == 404 or "not found" in lower:
                raise RuntimeError(
                    f"Ollama does not have model '{model_id}' installed. "
                    f"Run: ollama pull {model_id}"
                )
            raise RuntimeError(
                f"Ollama error for {model_name} ({model_id}): "
                f"HTTP {response.status_code} — {body}"
            )

        result = response.json()
        outline = self._extract_json_outline(result["response"])
        print(f"✓ Generated outline with {len(outline.get('slides', []))} slides")
        return outline

    def _generate_outline_claude(self, content: str, num_slides: int) -> Dict:
        """Generate outline via Anthropic Claude API."""
        try:
            from anthropic import Anthropic
        except ImportError as e:
            raise ImportError(
                "anthropic package required for Claude. Install with: pip install anthropic"
            ) from e

        api_key = os.getenv("ANTHROPIC_API_KEY") or self.claude_config.get("api_key")
        if not api_key:
            raise ValueError(
                "ANTHROPIC_API_KEY environment variable is required for Claude model"
            )

        model_id = self.claude_config.get("model_id", "claude-sonnet-5")
        prompt = self._build_prompt(content, num_slides)

        print(f"Generating outline with Claude ({model_id})...")

        client = Anthropic(api_key=api_key)
        message = client.messages.create(
            model=model_id,
            max_tokens=4000,
            messages=[{"role": "user", "content": prompt}],
        )

        # Claude Sonnet 5 may return thinking blocks — extract text only
        response_text = ""
        for block in message.content:
            if hasattr(block, "text"):
                response_text += block.text

        outline = self._extract_json_outline(response_text)
        print(f"✓ Generated outline with {len(outline.get('slides', []))} slides")
        return outline

    def generate_outline(self, content: str, model_key: str, num_slides: int = 10) -> Dict:
        """Generate presentation outline using the selected model."""
        if model_key not in self.models:
            raise ValueError(f"Unknown model: {model_key}. Available: {list(self.models)}")

        provider = self.models[model_key].get("provider", "ollama")
        if provider == "claude":
            return self._generate_outline_claude(content, num_slides)
        return self._generate_outline_ollama(content, model_key, num_slides)

    def create_presentation(
        self,
        outline: Dict,
        theme_key: str = "modern",
        output_path: str = "presentation.pptx",
        slide_transition: str = "fade",
        bullet_animation: str = "fade_in",
    ) -> str:
        """Create a polished PowerPoint presentation from an outline."""
        if theme_key not in self.themes:
            raise ValueError(f"Unknown theme: {theme_key}. Available: {list(self.themes)}")

        anim_cfg = self.config.get("animations", {})
        transitions = anim_cfg.get("slide_transitions", {"fade": "Fade", "none": "None"})
        bullet_anims = anim_cfg.get("bullet_animations", {"appear": "Appear", "none": "None"})
        if slide_transition not in transitions:
            slide_transition = "fade"
        if bullet_animation not in bullet_anims:
            bullet_animation = "fade_in"

        theme = self.themes[theme_key]
        self._active_theme = theme
        self._slide_transition = slide_transition
        self._bullet_animation = bullet_animation

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        primary = RGBColor(*theme["primary_color"])
        accent = RGBColor(*theme["accent_color"])
        bg = RGBColor(*theme["background"])
        text = RGBColor(*theme["text_color"])

        if "gradient_colors" in theme:
            grad_start = RGBColor(*theme["gradient_colors"][0])
            grad_end = RGBColor(*theme["gradient_colors"][1])
        else:
            grad_start = primary
            grad_end = accent

        print(
            f"Creating polished presentation with '{theme['name']}' "
            f"(transition={slide_transition}, bullets={bullet_animation})..."
        )

        for i, slide_data in enumerate(outline.get("slides", [])):
            slide_type = slide_data.get("type", "content")

            if slide_type == "title":
                slide = self._create_title_slide(
                    prs, slide_data, primary, accent, grad_start, grad_end
                )
            elif slide_type == "section":
                slide = self._create_section_slide(
                    prs, slide_data, primary, accent, grad_start, grad_end
                )
            elif slide_type == "conclusion":
                slide = self._create_conclusion_slide(prs, slide_data, primary, accent)
            else:
                slide = self._create_content_slide(prs, slide_data, primary, accent, text, bg)

            self._add_transition(slide, slide_transition)
            print(f"  ✓ Created slide {i + 1}: {slide_data.get('title', 'Untitled')}")

        prs.save(output_path)
        print(f"\n✓ Presentation saved: {output_path}")
        return output_path

    def _theme_image_path(self) -> Optional[str]:
        theme = getattr(self, "_active_theme", {}) or {}
        rel = theme.get("background_image")
        if not rel:
            return None
        path = rel if os.path.isabs(rel) else os.path.join(os.getcwd(), rel)
        return path if os.path.exists(path) else None

    def _apply_photo_background(self, slide, prs):
        """Add full-bleed photo background when theme provides an image."""
        path = self._theme_image_path()
        if not path:
            return False
        slide.shapes.add_picture(path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        return True

    def _add_scrim(self, slide, prs, color: RGBColor, transparency: float = 0.35):
        """Semi-transparent overlay so text stays readable on photo backgrounds."""
        scrim = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        scrim.fill.solid()
        scrim.fill.fore_color.rgb = color
        scrim.fill.transparency = transparency
        scrim.line.fill.background()
        return scrim

    def _add_transition(self, slide, transition: str = "fade"):
        """Add a PowerPoint slide transition via OOXML."""
        if not transition or transition == "none":
            return
        try:
            sld = slide._element
            for existing in sld.findall(f"{_P}transition"):
                sld.remove(existing)

            # Slow fade reads better; advance on click
            trans = etree.Element(f"{_P}transition", spd="slow", advClick="1")
            child_tag = {
                "fade": "fade",
                "push": "push",
                "wipe": "wipe",
            }.get(transition, "fade")
            if child_tag == "push":
                etree.SubElement(trans, f"{_P}push", dir="l")
            elif child_tag == "wipe":
                etree.SubElement(trans, f"{_P}wipe", dir="l")
            else:
                etree.SubElement(trans, f"{_P}fade")

            timing = sld.find(f"{_P}timing")
            if timing is not None:
                timing.addprevious(trans)
            else:
                c_sld = sld.find(f"{_P}cSld")
                if c_sld is not None:
                    c_sld.addnext(trans)
                else:
                    sld.append(trans)
        except Exception:
            pass

    def _shape_spid(self, shape) -> str:
        """Return the DrawingML shape id used by animation targets."""
        try:
            cNvPr = shape._element.nvSpPr.cNvPr
            return str(cNvPr.get("id"))
        except Exception:
            return str(shape.shape_id)

    def _anim_set_visibility(self, parent, shape_id: str, visible: bool, next_id_fn, delay: str = "0"):
        """OOXML <p:set> for style.visibility."""
        set_el = etree.SubElement(parent, f"{_P}set")
        c_bhvr = etree.SubElement(set_el, f"{_P}cBhvr")
        c_tn = etree.SubElement(
            c_bhvr,
            f"{_P}cTn",
            id=next_id_fn(),
            dur="1",
            fill="hold",
        )
        st = etree.SubElement(c_tn, f"{_P}stCondLst")
        etree.SubElement(st, f"{_P}cond", delay=delay)
        tgt = etree.SubElement(c_bhvr, f"{_P}tgtEl")
        etree.SubElement(tgt, f"{_P}spTgt", spid=shape_id)
        attr = etree.SubElement(c_bhvr, f"{_P}attrNameLst")
        etree.SubElement(attr, f"{_P}attrName").text = "style.visibility"
        to = etree.SubElement(set_el, f"{_P}to")
        etree.SubElement(to, f"{_P}strVal", val="visible" if visible else "hidden")

    def _anim_effect(
        self,
        parent,
        shape_id: str,
        next_id_fn,
        *,
        transition: str = "in",
        effect_filter: str = "fade",
        dur_ms: str = "500",
    ):
        """OOXML entrance/exit effect (matches PowerPoint-exported structure)."""
        anim = etree.SubElement(
            parent,
            f"{_P}animEffect",
            transition=transition,
            filter=effect_filter,
        )
        c_bhvr = etree.SubElement(anim, f"{_P}cBhvr")
        # Native PowerPoint omits fill on animEffect; holding here can leave text invisible.
        etree.SubElement(c_bhvr, f"{_P}cTn", id=next_id_fn(), dur=dur_ms)
        tgt = etree.SubElement(c_bhvr, f"{_P}tgtEl")
        etree.SubElement(tgt, f"{_P}spTgt", spid=shape_id)

    def _add_appear_animations(self, slide, shapes: List):
        """
        Sequential on-click text-line animations.

        Structure mirrors PowerPoint's own exports:
        clickEffect cTn → childTnLst → [set visibility, animEffect]
        (no extra nested par wrappers, which Mac PowerPoint often ignores).
        """
        style = getattr(self, "_bullet_animation", "fade_in")
        # Back-compat with older saved UI values
        if style == "fade":
            style = "fade_in"
        if not shapes or style == "none":
            return

        effect_filter = {
            "fade_in": "fade",
            "fade_in_out": "fade",
            "appear": "fade",
            "fly_left": "fly(fromLeft)",
        }.get(style, "fade")
        # PowerPoint preset IDs: 1=Appear, 10=Fade, 2=Fly
        preset_id = {
            "appear": "1",
            "fade_in": "10",
            "fade_in_out": "10",
            "fly_left": "2",
        }.get(style, "10")
        preset_subtype = "8" if style == "fly_left" else "0"  # 8 = from left
        enter_dur = "350" if style == "appear" else "500"
        do_fade_out = style == "fade_in_out"

        try:
            sld = slide._element
            for existing in sld.findall(f"{_P}timing"):
                sld.remove(existing)

            timing = etree.SubElement(sld, f"{_P}timing")
            tn_lst = etree.SubElement(timing, f"{_P}tnLst")
            root_par = etree.SubElement(tn_lst, f"{_P}par")
            root_tn = etree.SubElement(
                root_par,
                f"{_P}cTn",
                id="1",
                dur="indefinite",
                restart="never",
                nodeType="tmRoot",
            )
            root_children = etree.SubElement(root_tn, f"{_P}childTnLst")

            seq = etree.SubElement(root_children, f"{_P}seq", concurrent="1", nextAc="seek")
            seq_tn = etree.SubElement(
                seq,
                f"{_P}cTn",
                id="2",
                dur="indefinite",
                nodeType="mainSeq",
            )
            seq_children = etree.SubElement(seq_tn, f"{_P}childTnLst")

            next_id = 3

            def alloc_id() -> str:
                nonlocal next_id
                cur = str(next_id)
                next_id += 1
                return cur

            shape_ids = [self._shape_spid(shape) for shape in shapes]

            for index, spid in enumerate(shape_ids):
                step = etree.SubElement(seq_children, f"{_P}par")
                # grpId is required for PowerPoint to treat this as a real build effect
                step_tn = etree.SubElement(
                    step,
                    f"{_P}cTn",
                    id=alloc_id(),
                    fill="hold",
                    nodeType="clickEffect",
                    presetID=preset_id,
                    presetClass="entr",
                    presetSubtype=preset_subtype,
                    grpId=str(index),
                )
                step_st = etree.SubElement(step_tn, f"{_P}stCondLst")
                etree.SubElement(step_st, f"{_P}cond", delay="indefinite")
                # Direct children — same shape PowerPoint writes (no nested par/cTn)
                step_kids = etree.SubElement(step_tn, f"{_P}childTnLst")

                if do_fade_out and index > 0:
                    prev_id = shape_ids[index - 1]
                    self._anim_effect(
                        step_kids,
                        prev_id,
                        alloc_id,
                        transition="out",
                        effect_filter="fade",
                        dur_ms="400",
                    )
                    self._anim_set_visibility(step_kids, prev_id, False, alloc_id, delay="400")

                # Appear = visibility only; fade/fly also run animEffect
                self._anim_set_visibility(step_kids, spid, True, alloc_id, delay="0")
                if style != "appear":
                    self._anim_effect(
                        step_kids,
                        spid,
                        alloc_id,
                        transition="in",
                        effect_filter=effect_filter,
                        dur_ms=enter_dur,
                    )

            prev = etree.SubElement(seq, f"{_P}prevCondLst")
            etree.SubElement(prev, f"{_P}cond", evt="onPrev", delay="0")
            nxt = etree.SubElement(seq, f"{_P}nextCondLst")
            etree.SubElement(nxt, f"{_P}cond", evt="onNext", delay="0")
        except Exception as exc:
            print(f"Warning: could not add text-line animations: {exc}")

    def _create_title_slide(self, prs, data, primary, accent, grad_start, grad_end):
        """Polished title slide with geometric accents and optional photo background."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shapes = slide.shapes

        if self._apply_photo_background(slide, prs):
            self._add_scrim(slide, prs, primary, transparency=0.45)
        else:
            bg_shape1 = shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
            bg_shape1.fill.solid()
            bg_shape1.fill.fore_color.rgb = grad_start
            bg_shape1.line.fill.background()

            bg_shape2 = shapes.add_shape(1, 0, Inches(3), prs.slide_width, Inches(4.5))
            bg_shape2.fill.solid()
            bg_shape2.fill.fore_color.rgb = grad_end
            bg_shape2.fill.transparency = 0.3
            bg_shape2.line.fill.background()

        circle1 = shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(8), Inches(-0.5), Inches(2.5), Inches(2.5)
        )
        circle1.fill.solid()
        circle1.fill.fore_color.rgb = accent
        circle1.fill.transparency = 0.7
        circle1.line.fill.background()

        for i, pos in enumerate([(0.3, 6.8), (0.8, 6.5), (1.3, 6.9)]):
            square = shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(pos[0]),
                Inches(pos[1]),
                Inches(0.25),
                Inches(0.25),
            )
            square.fill.solid()
            square.fill.fore_color.rgb = accent
            square.fill.transparency = 0.3 * (i + 1)
            square.line.fill.background()

        title_box = shapes.add_textbox(Inches(1), Inches(2.3), Inches(8), Inches(1.8))
        title_frame = title_box.text_frame
        title_frame.text = data.get("title", "Presentation Title")
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        title_para.line_spacing = 1.1

        if data.get("subtitle"):
            subtitle_box = shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(7), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = data["subtitle"]
            subtitle_frame.word_wrap = True
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
            subtitle_para.alignment = PP_ALIGN.CENTER

        return slide

    def _create_section_slide(self, prs, data, primary, accent, grad_start, grad_end):
        """Polished section divider with geometric accents / photo background."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shapes = slide.shapes

        if self._apply_photo_background(slide, prs):
            self._add_scrim(slide, prs, primary, transparency=0.5)
        else:
            bg1 = shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
            bg1.fill.solid()
            bg1.fill.fore_color.rgb = grad_start
            bg1.line.fill.background()

            bg2 = shapes.add_shape(1, 0, Inches(2), prs.slide_width, Inches(5.5))
            bg2.fill.solid()
            bg2.fill.fore_color.rgb = grad_end
            bg2.fill.transparency = 0.4
            bg2.line.fill.background()

        for i, offset in enumerate([0, 0.3, 0.6]):
            bar = shapes.add_shape(
                1, Inches(-1 + offset), Inches(2.5 + i * 0.1), Inches(12), Inches(0.15)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent
            bar.fill.transparency = 0.4 + (i * 0.2)
            bar.line.fill.background()
            bar.rotation = 2

        circle = shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(5), Inches(3), Inches(3))
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.fill.transparency = 0.85
        circle.line.fill.background()

        title_box = shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(7), Inches(1.2))
        title_frame = title_box.text_frame
        title_frame.text = data.get("title", "Section")
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(52)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

        return slide

    def _create_conclusion_slide(self, prs, data, primary, accent):
        """Polished conclusion slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shapes = slide.shapes

        if self._apply_photo_background(slide, prs):
            self._add_scrim(slide, prs, accent, transparency=0.35)
        else:
            bg = shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = accent
            bg.line.fill.background()

        for i, size in enumerate([2.5, 2, 1.5]):
            ring = shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(3.75 + i * 0.25),
                Inches(2.5 + i * 0.25),
                Inches(size),
                Inches(size),
            )
            ring.fill.solid()
            ring.fill.fore_color.rgb = RGBColor(255, 255, 255)
            ring.fill.transparency = 0.8 + (i * 0.05)
            ring.line.fill.background()

        title_box = shapes.add_textbox(Inches(1.5), Inches(3.3), Inches(7), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = data.get("title", "Thank You")
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(56)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

        return slide

    def _create_content_slide(self, prs, data, primary, accent, text, bg):
        """Polished content slide: tight header, spaced bullets, animations."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shapes = slide.shapes
        is_photo = self._apply_photo_background(slide, prs)

        if is_photo:
            self._add_scrim(slide, prs, primary, transparency=0.45)
            panel = shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.55),
                Inches(0.45),
                Inches(8.9),
                Inches(6.6),
            )
            panel.fill.solid()
            panel.fill.fore_color.rgb = RGBColor(252, 252, 252)
            panel.fill.transparency = 0.06
            panel.line.fill.background()
            body_text = RGBColor(30, 30, 30)
        else:
            body_text = text

        header1 = shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.2))
        header1.fill.solid()
        header1.fill.fore_color.rgb = primary
        header1.line.fill.background()

        header2 = shapes.add_shape(1, 0, Inches(0.6), prs.slide_width, Inches(0.6))
        header2.fill.solid()
        header2.fill.fore_color.rgb = accent
        header2.fill.transparency = 0.7
        header2.line.fill.background()

        side_bar = shapes.add_shape(1, 0, Inches(1.2), Inches(0.2), Inches(6.3))
        side_bar.fill.solid()
        side_bar.fill.fore_color.rgb = accent
        side_bar.fill.transparency = 0.3
        side_bar.line.fill.background()

        for pos in [(0.35, 0.5), (0.6, 0.5)]:
            dot = shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(pos[0]), Inches(pos[1]), Inches(0.12), Inches(0.12)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = accent
            dot.line.fill.background()

        title_box = shapes.add_textbox(Inches(0.9), Inches(0.25), Inches(8.5), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = data.get("title", "Slide Title")
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        underline = shapes.add_shape(1, Inches(0.9), Inches(1.2), Inches(1.8), Inches(0.08))
        underline.fill.solid()
        underline.fill.fore_color.rgb = accent
        underline.line.fill.background()

        bullets = (data.get("bullets") or [])[: self.MAX_BULLETS]
        animated_shapes = []
        if bullets:
            start_y = 2.0
            spacing = 0.95
            for i, bullet in enumerate(bullets):
                y_pos = start_y + (i * spacing)
                bullet_box = shapes.add_textbox(
                    Inches(1.2), Inches(y_pos), Inches(7.5), Inches(0.75)
                )
                text_frame = bullet_box.text_frame
                text_frame.word_wrap = True
                p = text_frame.paragraphs[0]
                p.text = "• " + bullet
                p.font.size = Pt(20)
                p.font.color.rgb = body_text
                p.line_spacing = 1.2
                animated_shapes.append(bullet_box)

        corner = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9), Inches(6.9), Inches(0.7), Inches(0.4)
        )
        corner.fill.solid()
        corner.fill.fore_color.rgb = accent
        corner.fill.transparency = 0.7
        corner.line.fill.background()

        self._add_appear_animations(slide, animated_shapes)
        return slide

    def generate_from_text(
        self,
        content: str,
        model_key: str,
        num_slides: int = 10,
        theme: str = "modern",
        output_path: str = "presentation.pptx",
        slide_transition: str = "fade",
        bullet_animation: str = "fade_in",
    ) -> str:
        """Complete pipeline: text -> outline -> polished designer presentation."""
        print(f"\n{'=' * 60}")
        print("AI Presentation Generator (Polished Designer)")
        print(f"{'=' * 60}")
        print(f"Model: {self.models[model_key]['name']}")
        print(f"Theme: {self.themes[theme]['name']}")
        print(f"Slides: {num_slides}")
        print(f"Transition: {slide_transition} | Bullet animation: {bullet_animation}")
        print(f"{'=' * 60}\n")

        outline = self.generate_outline(content, model_key, num_slides)
        result_path = self.create_presentation(
            outline,
            theme,
            output_path,
            slide_transition=slide_transition,
            bullet_animation=bullet_animation,
        )

        print(f"\n{'=' * 60}")
        print("✓ SUCCESS! Presentation created")
        print(f"{'=' * 60}\n")

        return result_path


if __name__ == "__main__":
    generator = PresentationGenerator()

    print("\nAvailable Models:")
    for key, model in generator.list_models().items():
        print(f"  - {key}: {model['name']} - {model['description']}")

    print("\nAvailable Themes:")
    for key, theme in generator.list_themes().items():
        print(f"  - {key}: {theme['name']}")
