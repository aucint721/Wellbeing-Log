#!/usr/bin/env python3
"""
Real Claude API Presentation Generator Demo
Shows TRUE quality with actual Claude Opus 5 Sonnet
"""

import os
import json
import yaml
from typing import Dict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def get_claude_outline(api_key: str, topic: str, num_slides: int = 10) -> Dict:
    """Generate presentation outline using real Claude API"""
    
    try:
        from anthropic import Anthropic
    except ImportError:
        print("Installing anthropic package...")
        import subprocess
        subprocess.check_call(["pip", "install", "--user", "anthropic"])
        from anthropic import Anthropic
    
    client = Anthropic(api_key=api_key)
    
    prompt = f"""You are an expert presentation designer. Create a compelling, professional presentation about: {topic}

Requirements:
- Create exactly {num_slides} slides
- Make the content engaging, insightful, and well-structured
- Each content slide should have 3-5 concise, impactful bullet points
- Use clear, professional language
- Include relevant examples and insights
- Organize with logical flow and section breaks

Return ONLY valid JSON in this exact format:
{{
  "title": "Main presentation title",
  "slides": [
    {{
      "type": "title",
      "title": "Title text",
      "subtitle": "Subtitle text"
    }},
    {{
      "type": "content",
      "title": "Slide title",
      "bullets": ["Point 1", "Point 2", "Point 3", "Point 4", "Point 5"]
    }},
    {{
      "type": "section",
      "title": "Section divider title"
    }},
    {{
      "type": "conclusion",
      "title": "Thank You"
    }}
  ]
}}

Slide types:
- "title": Opening slide with title and subtitle
- "content": Regular slide with bullets (3-5 bullets each)
- "section": Section divider with just a title (use sparingly for major transitions)
- "conclusion": Final slide

Create an engaging, professional presentation now:"""

    print(f"🤖 Asking Claude to create presentation about: {topic}")
    print("⏳ This may take 10-20 seconds...\n")
    
    message = client.messages.create(
        model="claude-sonnet-5",
        max_tokens=4000,
        messages=[
            {"role": "user", "content": prompt}
        ]
    )
    
    response_text = message.content[0].text.strip()
    
    # Extract JSON
    json_start = response_text.find('{')
    json_end = response_text.rfind('}') + 1
    
    if json_start >= 0 and json_end > json_start:
        json_text = response_text[json_start:json_end]
        outline = json.loads(json_text)
        print(f"✅ Claude generated: {len(outline.get('slides', []))} slides")
        return outline
    else:
        raise ValueError("No valid JSON found in Claude's response")


class ThemePreview:
    """Generate presentations with different themes"""
    
    def __init__(self):
        with open('config.yaml', 'r') as f:
            config = yaml.safe_load(f)
        self.themes = config['themes']
    
    def create_presentation(self, outline: Dict, theme_name: str, output_path: str):
        """Create PowerPoint presentation with specified theme"""
        
        theme = self.themes[theme_name]
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Color scheme
        primary = RGBColor(*theme['primary_color'])
        accent = RGBColor(*theme['accent_color'])
        bg = RGBColor(*theme['background'])
        text = RGBColor(*theme['text_color'])
        
        print(f"\n🎨 Creating '{theme['name']}' presentation...")
        
        for i, slide_data in enumerate(outline['slides']):
            slide_type = slide_data.get('type', 'content')
            
            if slide_type == 'title':
                self._create_title_slide(prs, slide_data, primary, accent, text)
            elif slide_type == 'section':
                self._create_section_slide(prs, slide_data, primary, text)
            elif slide_type == 'conclusion':
                self._create_conclusion_slide(prs, slide_data, accent, text)
            else:
                self._create_content_slide(prs, slide_data, primary, accent, text)
            
            print(f"  ✓ Slide {i+1}/{len(outline['slides'])}: {slide_data.get('title', 'Untitled')}")
        
        prs.save(output_path)
        print(f"✅ Saved: {output_path}")
        return output_path
    
    def _create_title_slide(self, prs, data, primary_color, accent_color, text_color):
        """Create title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        shapes = slide.shapes
        bg_shape = shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = primary_color
        bg_shape.line.fill.background()
        
        # Accent bar at bottom
        accent_bar = shapes.add_shape(1, 0, Inches(6.5), prs.slide_width, Inches(1))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_color
        accent_bar.line.fill.background()
        
        # Title
        title_box = shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = data.get('title', 'Presentation Title')
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if data.get('subtitle'):
            subtitle_box = shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = data['subtitle']
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(28)
            subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
            subtitle_para.alignment = PP_ALIGN.CENTER
    
    def _create_section_slide(self, prs, data, primary_color, text_color):
        """Create section divider slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        shapes = slide.shapes
        bg_shape = shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = primary_color
        bg_shape.line.fill.background()
        
        title_box = shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = data.get('title', 'Section')
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
    
    def _create_conclusion_slide(self, prs, data, accent_color, text_color):
        """Create conclusion slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        shapes = slide.shapes
        bg_shape = shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = accent_color
        bg_shape.line.fill.background()
        
        title_box = shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = data.get('title', 'Thank You')
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
    
    def _create_content_slide(self, prs, data, primary_color, accent_color, text_color):
        """Create content slide with bullets"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shapes = slide.shapes
        
        # Title bar
        title_shape = shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.2))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = primary_color
        title_shape.line.fill.background()
        
        # Title text
        title_box = shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = data.get('title', 'Slide Title')
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Accent line
        accent_line = shapes.add_shape(1, 0, Inches(1.2), prs.slide_width, Inches(0.05))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = accent_color
        accent_line.line.fill.background()
        
        # Content bullets
        bullets = data.get('bullets', [])
        if bullets:
            content_box = shapes.add_textbox(Inches(1), Inches(1.8), Inches(8.5), Inches(5.2))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = bullet
                p.level = 0
                p.font.size = Pt(18)
                p.font.color.rgb = text_color
                p.space_after = Pt(16)
                p.line_spacing = 1.2


def main():
    """Generate presentations with real Claude API"""
    
    print("\n" + "="*70)
    print("🤖 REAL CLAUDE API PRESENTATION GENERATOR")
    print("="*70)
    
    # Get API key
    api_key = os.getenv('ANTHROPIC_API_KEY')
    if not api_key:
        api_key = input("\n🔑 Enter your Claude API key: ").strip()
    
    if not api_key:
        print("❌ API key required!")
        return
    
    # Get topic
    print("\n" + "="*70)
    topic = input("📝 What should the presentation be about?\n   (e.g., 'Machine Learning for Beginners', 'Climate Change Solutions')\n\n   Topic: ").strip()
    
    if not topic:
        topic = "The Future of Artificial Intelligence"
        print(f"\n   Using default: {topic}")
    
    # Get number of slides
    num_input = input("\n📊 How many slides? (default: 12): ").strip()
    num_slides = int(num_input) if num_input.isdigit() else 12
    
    # Get theme preference
    print("\n🎨 Available themes:")
    print("   1. modern - Professional blue & orange (default)")
    print("   2. dark - Tech-style dark theme")
    print("   3. education - Friendly educational theme")
    print("   4. warm - Creative warm colors")
    print("   5. all - Generate all 4 themes")
    
    theme_input = input("\n   Theme choice (1-5): ").strip()
    theme_map = {
        '1': ['modern'],
        '2': ['dark'],
        '3': ['education'],
        '4': ['warm'],
        '5': ['modern', 'dark', 'education', 'warm']
    }
    themes = theme_map.get(theme_input, ['modern'])
    
    print("\n" + "="*70)
    print("🚀 GENERATING PRESENTATION")
    print("="*70)
    
    # Generate outline with Claude
    outline = get_claude_outline(api_key, topic, num_slides)
    
    print(f"\n📊 Presentation: {outline['title']}")
    print(f"📝 Slides: {len(outline['slides'])}")
    
    # Create presentations
    previewer = ThemePreview()
    
    for theme_name in themes:
        output = f"claude_{theme_name}_theme.pptx"
        previewer.create_presentation(outline, theme_name, output)
    
    print("\n" + "="*70)
    print("✅ DONE! Real Claude-powered presentations created!")
    print("="*70)
    print(f"""
Files created:
""")
    for theme_name in themes:
        print(f"- claude_{theme_name}_theme.pptx")
    
    print(f"""
🎯 This is the REAL QUALITY you'll get!
   - Content by Claude Opus 5 Sonnet
   - Professional slide design
   - Fully editable PowerPoint
   
Open the files to see the difference! 🎨
""")


if __name__ == "__main__":
    main()
