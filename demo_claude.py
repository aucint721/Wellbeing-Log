#!/usr/bin/env python3
"""
Demo: Claude Opus 5 Presentation Generator
Test the quality before building the full system
"""

import json
import yaml
from typing import Dict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# DEMO: This simulates what Claude Opus 5 would return
# In the real version, this comes from Claude API
def get_claude_demo_response():
    """Simulated Claude Opus 5 response for demo purposes"""
    return {
        "title": "Introduction to Artificial Intelligence",
        "slides": [
            {
                "type": "title",
                "title": "Introduction to Artificial Intelligence",
                "subtitle": "Understanding the Future of Technology"
            },
            {
                "type": "content",
                "title": "What is Artificial Intelligence?",
                "bullets": [
                    "AI enables machines to perform tasks requiring human intelligence",
                    "Includes learning, reasoning, problem-solving, and perception",
                    "Narrow AI: Specific tasks (like Siri, recommendation systems)",
                    "General AI: Human-level intelligence across domains (future goal)",
                    "AI has evolved from 1950s theoretical concepts to today's practical applications"
                ]
            },
            {
                "type": "section",
                "title": "Core Technologies"
            },
            {
                "type": "content",
                "title": "Machine Learning Fundamentals",
                "bullets": [
                    "Supervised Learning: Learning from labeled examples (classification, regression)",
                    "Unsupervised Learning: Finding patterns in unlabeled data (clustering)",
                    "Reinforcement Learning: Learning through trial and error with rewards",
                    "Deep Learning: Neural networks with multiple layers for complex patterns",
                    "Models improve accuracy through iterative training on large datasets"
                ]
            },
            {
                "type": "content",
                "title": "Neural Networks & Deep Learning",
                "bullets": [
                    "Inspired by biological neurons in the human brain",
                    "Layers: Input → Hidden layers (feature extraction) → Output (prediction)",
                    "Activation functions introduce non-linearity for complex problem-solving",
                    "Computer Vision: Image recognition, object detection, facial recognition",
                    "Natural Language Processing: Translation, sentiment analysis, text generation"
                ]
            },
            {
                "type": "section",
                "title": "Real-World Impact"
            },
            {
                "type": "content",
                "title": "Applications Across Industries",
                "bullets": [
                    "Healthcare: Early disease detection, drug discovery, personalized medicine",
                    "Finance: Fraud detection, algorithmic trading, risk assessment",
                    "Transportation: Autonomous vehicles, route optimization, traffic prediction",
                    "Retail: Recommendation engines, inventory management, customer insights",
                    "Entertainment: Content recommendations, game AI, music composition"
                ]
            },
            {
                "type": "content",
                "title": "Ethical Considerations",
                "bullets": [
                    "Bias & Fairness: AI systems can inherit biases from training data",
                    "Privacy: Data collection and surveillance concerns require careful governance",
                    "Transparency: 'Black box' models need explainability for trust",
                    "Job Displacement: Automation may transform workforce, requiring reskilling",
                    "Safety & Alignment: Ensuring AI systems act in accordance with human values"
                ]
            },
            {
                "type": "content",
                "title": "The Future of AI",
                "bullets": [
                    "Multimodal AI: Systems that understand text, images, audio, and video together",
                    "Edge AI: Processing on devices for faster, more private applications",
                    "AI Safety Research: Developing robust, controllable, and beneficial AI",
                    "Human-AI Collaboration: Augmenting rather than replacing human capabilities",
                    "Career Opportunities: Data scientists, ML engineers, AI researchers in high demand"
                ]
            },
            {
                "type": "conclusion",
                "title": "Thank You"
            }
        ]
    }


class ThemePreview:
    """Generate preview presentations with different themes"""
    
    def __init__(self):
        # Load themes from config
        with open('config.yaml', 'r') as f:
            config = yaml.safe_load(f)
        self.themes = config['themes']
    
    def create_presentation(self, outline: Dict, theme_name: str, output_path: str):
        """Create PowerPoint presentation with specified theme"""
        
        theme = self.themes[theme_name]
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Color scheme
        primary = RGBColor(*theme['primary_color'])
        accent = RGBColor(*theme['accent_color'])
        bg = RGBColor(*theme['background'])
        text = RGBColor(*theme['text_color'])
        
        print(f"\n🎨 Creating presentation with '{theme['name']}' theme...")
        
        for i, slide_data in enumerate(outline['slides']):
            slide_type = slide_data.get('type', 'content')
            
            if slide_type == 'title':
                self._create_title_slide(prs, slide_data, primary, accent, text)
            elif slide_type == 'section':
                self._create_section_slide(prs, slide_data, primary, text)
            elif slide_type == 'conclusion':
                self._create_conclusion_slide(prs, slide_data, accent, text)
            else:
                self._create_content_slide(prs, slide_data, primary, accent, text)
            
            print(f"  ✓ Slide {i+1}/{len(outline['slides'])}: {slide_data.get('title', 'Untitled')}")
        
        prs.save(output_path)
        print(f"\n✅ Presentation saved: {output_path}")
        print(f"📊 Total slides: {len(outline['slides'])}")
        return output_path
    
    def _create_title_slide(self, prs, data, primary_color, accent_color, text_color):
        """Create title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background rectangle
        shapes = slide.shapes
        bg_shape = shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = primary_color
        bg_shape.line.fill.background()
        
        # Accent bar at bottom
        accent_bar = shapes.add_shape(1, 0, Inches(6.5), prs.slide_width, Inches(1))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_color
        accent_bar.line.fill.background()
        
        # Title
        title_box = shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = data.get('title', 'Presentation Title')
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if data.get('subtitle'):
            subtitle_box = shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = data['subtitle']
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(28)
            subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
            subtitle_para.alignment = PP_ALIGN.CENTER
    
    def _create_section_slide(self, prs, data, primary_color, text_color):
        """Create section divider slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        shapes = slide.shapes
        bg_shape = shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = primary_color
        bg_shape.line.fill.background()
        
        # Section title
        title_box = shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = data.get('title', 'Section')
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
    
    def _create_conclusion_slide(self, prs, data, accent_color, text_color):
        """Create conclusion slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        shapes = slide.shapes
        bg_shape = shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = accent_color
        bg_shape.line.fill.background()
        
        title_box = shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = data.get('title', 'Thank You')
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
    
    def _create_content_slide(self, prs, data, primary_color, accent_color, text_color):
        """Create content slide with bullets"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shapes = slide.shapes
        
        # Title bar with gradient effect (simulated with two shapes)
        title_shape = shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.2))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = primary_color
        title_shape.line.fill.background()
        
        # Title text
        title_box = shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = data.get('title', 'Slide Title')
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Accent line under header
        accent_line = shapes.add_shape(1, 0, Inches(1.2), prs.slide_width, Inches(0.05))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = accent_color
        accent_line.line.fill.background()
        
        # Content bullets
        bullets = data.get('bullets', [])
        if bullets:
            content_box = shapes.add_textbox(Inches(1), Inches(1.8), Inches(8.5), Inches(5.2))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = bullet
                p.level = 0
                p.font.size = Pt(18)
                p.font.color.rgb = text_color
                p.space_after = Pt(16)
                p.line_spacing = 1.2
    
    def generate_all_theme_demos(self, outline: Dict):
        """Generate demo presentations for all themes"""
        print("\n" + "="*60)
        print("🎨 GENERATING THEME DEMOS")
        print("="*60)
        
        for theme_name, theme_data in self.themes.items():
            output = f"demo_{theme_name}_theme.pptx"
            print(f"\n📌 Theme: {theme_data['name']}")
            self.create_presentation(outline, theme_name, output)
        
        print("\n" + "="*60)
        print("✅ ALL DEMOS GENERATED!")
        print("="*60)
        print("\nOpen the .pptx files to see the quality!")


def main():
    """Demo the presentation quality"""
    
    print("\n" + "="*70)
    print("🎨 CLAUDE OPUS 5 PRESENTATION QUALITY DEMO")
    print("="*70)
    print("\nThis demo shows you EXACTLY what quality you'll get!")
    print("Using simulated Claude Opus 5 output for AI content...")
    
    # Get demo content (this would come from Claude API in real version)
    outline = get_claude_demo_response()
    
    print(f"\n📊 Demo Presentation: {outline['title']}")
    print(f"📝 Slides: {len(outline['slides'])}")
    
    # Create theme previewer
    previewer = ThemePreview()
    
    # Generate demos for all themes
    previewer.generate_all_theme_demos(outline)
    
    print("\n" + "="*70)
    print("🎯 NEXT STEPS:")
    print("="*70)
    print("""
1. Open the generated .pptx files to see the quality
2. Check different themes: modern, dark, education, warm
3. These show EXACT output you'll get with Claude Opus 5
4. Quality of AI content will be even BETTER with real Claude

Files created:
- demo_modern_theme.pptx
- demo_dark_theme.pptx
- demo_education_theme.pptx
- demo_warm_theme.pptx
    """)


if __name__ == "__main__":
    main()
