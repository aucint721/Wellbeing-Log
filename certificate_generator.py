#!/usr/bin/env python3
"""
Certificate Generator
Creates professional certificate slides with logo placeholders.
Perfect for recognizing teacher aides and other school staff.
"""

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


class CertificateGenerator:
    """Generate certificate slides with customizable text and logo placeholders."""

    def __init__(self, config_path: str = "config.yaml"):
        with open(config_path, "r") as f:
            self.config = yaml.safe_load(f)
        self.themes = self.config["themes"]

    def create_certificate(
        self,
        output_file: str,
        recipient_name: str = "[Recipient Name]",
        title: str = "Certificate of Recognition",
        body_text: str = "In recognition of outstanding dedication and service as a Teacher Aide",
        date: str = "[Date]",
        signature_line: str = "[Principal/Administrator Name]",
        signature_title: str = "Principal",
        theme: str = "royal_purple"
    ):
        """
        Create a certificate presentation.
        
        Args:
            output_file: Output .pptx filename
            recipient_name: Name of the recipient (or placeholder)
            title: Certificate title
            body_text: Main recognition text
            date: Date of recognition
            signature_line: Name for signature line
            signature_title: Title below signature (e.g., "Principal", "Administrator", leave blank for none)
            theme: Theme from config.yaml (default: royal_purple for formal certificates)
        """
        
        # Get theme colors
        if theme not in self.themes:
            theme = "royal_purple"  # Fallback to formal theme
        
        theme_config = self.themes[theme]
        primary_color = RGBColor(*theme_config["primary_color"])
        accent_color = RGBColor(*theme_config["accent_color"])
        bg_color = RGBColor(*theme_config["background"])
        text_color = RGBColor(*theme_config["text_color"])
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Add blank slide
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        # Add decorative border
        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(6.5)
        )
        border.fill.background()
        border.line.color.rgb = primary_color
        border.line.width = Pt(8)
        
        # Inner decorative border
        inner_border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.75), Inches(0.75),
            Inches(8.5), Inches(6)
        )
        inner_border.fill.background()
        inner_border.line.color.rgb = accent_color
        inner_border.line.width = Pt(2)
        
        # Logo placeholder box (top center)
        logo_placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4.25), Inches(1.0),
            Inches(1.5), Inches(1.5)
        )
        logo_placeholder.fill.solid()
        logo_placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
        logo_placeholder.line.color.rgb = primary_color
        logo_placeholder.line.width = Pt(2)
        logo_placeholder.line.dash_style = 2  # Dashed line
        
        # Logo placeholder text
        logo_text = logo_placeholder.text_frame
        logo_text.text = "INSERT\nSCHOOL\nLOGO\nHERE"
        logo_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        for paragraph in logo_text.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.color.rgb = RGBColor(150, 150, 150)
            paragraph.font.bold = True
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.75),
            Inches(7), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = primary_color
        
        # Presented to
        presented_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(3.4),
            Inches(7), Inches(0.3)
        )
        presented_frame = presented_box.text_frame
        presented_frame.text = "Presented to"
        presented_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        presented_frame.paragraphs[0].font.size = Pt(16)
        presented_frame.paragraphs[0].font.italic = True
        presented_frame.paragraphs[0].font.color.rgb = text_color
        
        # Recipient name
        name_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(3.75),
            Inches(7), Inches(0.5)
        )
        name_frame = name_box.text_frame
        name_frame.text = recipient_name
        name_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        name_frame.paragraphs[0].font.size = Pt(32)
        name_frame.paragraphs[0].font.bold = True
        name_frame.paragraphs[0].font.color.rgb = accent_color
        
        # Body text
        body_box = slide.shapes.add_textbox(
            Inches(2), Inches(4.4),
            Inches(6), Inches(0.8)
        )
        body_frame = body_box.text_frame
        body_frame.text = body_text
        body_frame.word_wrap = True
        body_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        body_frame.paragraphs[0].font.size = Pt(16)
        body_frame.paragraphs[0].font.color.rgb = text_color
        
        # Date
        date_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(5.5),
            Inches(3), Inches(0.4)
        )
        date_frame = date_box.text_frame
        date_frame.text = date
        date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        date_frame.paragraphs[0].font.size = Pt(14)
        date_frame.paragraphs[0].font.color.rgb = text_color
        
        # Signature line
        sig_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5.5), Inches(5.65),
            Inches(3), Inches(0.01)
        )
        sig_line.fill.solid()
        sig_line.fill.fore_color.rgb = primary_color
        sig_line.line.fill.background()
        
        # Signature name
        sig_box = slide.shapes.add_textbox(
            Inches(5.5), Inches(5.75),
            Inches(3), Inches(0.4)
        )
        sig_frame = sig_box.text_frame
        sig_frame.text = signature_line
        sig_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        sig_frame.paragraphs[0].font.size = Pt(14)
        sig_frame.paragraphs[0].font.color.rgb = text_color
        
        # Signature title (only add if not empty)
        if signature_title:
            sig_title_box = slide.shapes.add_textbox(
                Inches(5.5), Inches(6.05),
                Inches(3), Inches(0.3)
            )
            sig_title_frame = sig_title_box.text_frame
            sig_title_frame.text = signature_title
            sig_title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            sig_title_frame.paragraphs[0].font.size = Pt(11)
            sig_title_frame.paragraphs[0].font.italic = True
            sig_title_frame.paragraphs[0].font.color.rgb = text_color
        
        # Save
        prs.save(output_file)
        print(f"✓ Certificate created: {output_file}")
        print(f"  Theme: {theme_config['name']}")
        print(f"  Recipient: {recipient_name}")
        print(f"\nTo add your school logo:")
        print(f"  1. Open {output_file} in PowerPoint or Google Slides")
        print(f"  2. Click on the 'INSERT SCHOOL LOGO HERE' box")
        print(f"  3. Delete it and insert your school logo image")
        print(f"  4. Position and resize the logo as needed")


def interactive_mode():
    """Interactive mode - ask user for all options."""
    import sys
    
    generator = CertificateGenerator()
    
    print("\n" + "="*60)
    print("         CERTIFICATE GENERATOR - INTERACTIVE MODE")
    print("="*60)
    print()
    
    # Output filename
    output_file = input("Output filename [teacher_aide_certificate.pptx]: ").strip()
    if not output_file:
        output_file = "teacher_aide_certificate.pptx"
    
    # Show theme options
    print("\n" + "-"*60)
    print("AVAILABLE THEMES:")
    print("-"*60)
    
    # Organize themes by category
    recommended = ["royal_purple", "education", "midnight_blue", "forest_minimal"]
    modern = ["sunset_gradient", "ocean_deep", "lavender_dream", "mint_fresh", "coral_pink"]
    photo = [k for k in generator.themes.keys() if k.startswith("photo_")]
    
    print("\n📌 RECOMMENDED FOR CERTIFICATES:")
    for i, key in enumerate(recommended, 1):
        if key in generator.themes:
            print(f"  {i}. {key:20} - {generator.themes[key]['name']}")
    
    print("\n🎨 MODERN THEMES:")
    start = len(recommended) + 1
    for i, key in enumerate(modern, start):
        if key in generator.themes:
            print(f"  {i}. {key:20} - {generator.themes[key]['name']}")
    
    print(f"\n📷 PHOTO BACKGROUNDS: (type 'photo' to see {len(photo)} photo themes)")
    
    print("-"*60)
    
    # Theme selection
    theme_input = input(f"\nSelect theme [1-{start + len(modern) - 1}] or name [royal_purple]: ").strip()
    
    if theme_input.lower() == 'photo':
        print("\n📷 PHOTO THEMES:")
        for i, key in enumerate(photo, 1):
            print(f"  {i}. {key:25} - {generator.themes[key]['name']}")
        photo_choice = input(f"\nSelect photo theme [1-{len(photo)}] or name: ").strip()
        if photo_choice.isdigit() and 1 <= int(photo_choice) <= len(photo):
            theme = photo[int(photo_choice) - 1]
        elif photo_choice in photo:
            theme = photo_choice
        else:
            theme = "royal_purple"
    elif theme_input.isdigit():
        choice_num = int(theme_input)
        all_themes = recommended + modern
        if 1 <= choice_num <= len(all_themes):
            theme = all_themes[choice_num - 1]
        else:
            theme = "royal_purple"
    elif theme_input in generator.themes:
        theme = theme_input
    else:
        theme = "royal_purple"
    
    print(f"✓ Selected theme: {generator.themes[theme]['name']}")
    
    # Certificate details
    print("\n" + "-"*60)
    print("CERTIFICATE DETAILS:")
    print("-"*60)
    
    title = input("\nCertificate title [Certificate of Recognition]: ").strip()
    if not title:
        title = "Certificate of Recognition"
    
    recipient_name = input("Recipient name [[Recipient Name] for blank template]: ").strip()
    if not recipient_name:
        recipient_name = "[Recipient Name]"
    
    body_text = input("Recognition text [In recognition of outstanding dedication...]: ").strip()
    if not body_text:
        body_text = "In recognition of outstanding dedication and service as a Teacher Aide"
    
    date = input("Date [[Date] for blank template]: ").strip()
    if not date:
        date = "[Date]"
    
    signature_line = input("Signature name [[Principal/Administrator Name]]: ").strip()
    if not signature_line:
        signature_line = "[Principal/Administrator Name]"
    
    signature_title = input("Title below signature (e.g., Principal, Administrator) [Principal]: ").strip()
    if signature_title == "":
        signature_title = "Principal"
    
    print("\n💡 Tip: Leave signature title blank (type 'none') if your signature name already includes the title")
    if signature_title.lower() in ['none', 'blank', 'empty', '']:
        signature_title = ""
    
    # Confirmation
    print("\n" + "="*60)
    print("PREVIEW:")
    print("="*60)
    print(f"Output file: {output_file}")
    print(f"Theme: {generator.themes[theme]['name']}")
    print(f"Title: {title}")
    print(f"Recipient: {recipient_name}")
    print(f"Body: {body_text}")
    print(f"Date: {date}")
    print(f"Signature: {signature_line}")
    print(f"Signature title: {signature_title if signature_title else '(none)'}")
    print("="*60)
    
    confirm = input("\nCreate certificate? [Y/n]: ").strip().lower()
    if confirm and confirm != 'y' and confirm != 'yes':
        print("Cancelled.")
        sys.exit(0)
    
    print()
    generator.create_certificate(
        output_file=output_file,
        recipient_name=recipient_name,
        title=title,
        body_text=body_text,
        date=date,
        signature_line=signature_line,
        signature_title=signature_title,
        theme=theme
    )


def main():
    """CLI interface for certificate generation."""
    import argparse
    import sys
    
    # Check if running with no arguments (interactive mode)
    if len(sys.argv) == 1:
        interactive_mode()
        return
    
    parser = argparse.ArgumentParser(
        description="Generate professional certificate slides for teacher aides"
    )
    parser.add_argument(
        "-o", "--output",
        default="teacher_aide_certificate.pptx",
        help="Output PowerPoint file (default: teacher_aide_certificate.pptx)"
    )
    parser.add_argument(
        "-n", "--name",
        default="[Recipient Name]",
        help="Recipient name (leave as placeholder for template)"
    )
    parser.add_argument(
        "-t", "--title",
        default="Certificate of Recognition",
        help="Certificate title"
    )
    parser.add_argument(
        "-b", "--body",
        default="In recognition of outstanding dedication and service as a Teacher Aide",
        help="Body text"
    )
    parser.add_argument(
        "-d", "--date",
        default="[Date]",
        help="Date of recognition"
    )
    parser.add_argument(
        "-s", "--signature",
        default="[Principal/Administrator Name]",
        help="Name for signature line"
    )
    parser.add_argument(
        "--signature-title",
        default="Principal",
        help="Title below signature (default: Principal). Use empty string for none."
    )
    parser.add_argument(
        "--theme",
        default="royal_purple",
        help="Theme name from config.yaml (default: royal_purple)"
    )
    parser.add_argument(
        "--list-themes",
        action="store_true",
        help="List available themes"
    )
    parser.add_argument(
        "-i", "--interactive",
        action="store_true",
        help="Run in interactive mode (asks for all options)"
    )
    
    args = parser.parse_args()
    
    generator = CertificateGenerator()
    
    if args.interactive:
        interactive_mode()
        return
    
    if args.list_themes:
        print("\nAvailable Certificate Themes:")
        print("-" * 50)
        for key, theme in generator.themes.items():
            print(f"  {key:20} - {theme['name']}")
        print("\nRecommended for certificates:")
        print("  royal_purple, education, midnight_blue, forest_minimal")
        return
    
    generator.create_certificate(
        output_file=args.output,
        recipient_name=args.name,
        title=args.title,
        body_text=args.body,
        date=args.date,
        signature_line=args.signature,
        signature_title=args.signature_title,
        theme=args.theme
    )


if __name__ == "__main__":
    main()
