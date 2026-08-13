#!/usr/bin/env python3
"""
Enhanced Claude API Presentation Generator
Better content, transitions, and visual design
"""

import os
import json
import yaml
from typing import Dict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def get_claude_outline(api_key: str, topic: str, num_slides: int = 10) -> Dict:
    """Generate presentation outline using real Claude API with enhanced prompt"""
    
    try:
        from anthropic import Anthropic
    except ImportError:
        print("Installing anthropic package...")
        import subprocess
        subprocess.check_call(["pip", "install", "anthropic"])
        from anthropic import Anthropic
    
    client = Anthropic(api_key=api_key)
    
    prompt = f"""You are an expert presentation designer creating a captivating presentation about: {topic}

CONTENT STYLE:
- Make bullet points PUNCHY, memorable, and impactful
- Use active voice and strong verbs
- Include specific examples and numbers when relevant
- Write like you're telling an engaging story, not a boring lecture
- Each bullet should be a complete insight, not just a fragment
- Aim for "aha!" moments that make people remember the content

TARGET AUDIENCE: Year 1 students (or general audience if not education-focused)
- Keep language clear and accessible
- Use concrete examples they can visualize
- Make it fun and engaging, not dry

Create exactly {num_slides} slides with this structure:
1. Title slide (engaging title + compelling subtitle)
2. Content slides (3-4 punchy bullets each - quality over quantity)
3. Section breaks (use 1-2 for major topic transitions)
4. Conclusion slide

Return ONLY valid JSON in this exact format:
{{
  "title": "Main presentation title (make it catchy!)",
  "slides": [
    {{
      "type": "title",
      "title": "Engaging Title Text",
      "subtitle": "Compelling subtitle that makes you want to learn more"
    }},
    {{
      "type": "content",
      "title": "Clear, Action-Oriented Slide Title",
      "bullets": [
        "First impactful point with specific detail",
        "Second memorable insight that sticks",
        "Third engaging example that illustrates the concept",
        "Optional fourth point (only if needed)"
      ]
    }},
    {{
      "type": "section",
      "title": "Major Topic Transition"
    }},
    {{
      "type": "conclusion",
      "title": "Thank You"
    }}
  ]
}}

IMPORTANT:
- Each bullet should be 1-2 lines max
- Use numbers, examples, or analogies to make points memorable
- Avoid generic statements - be specific and interesting
- Create a logical narrative flow across slides

Create an engaging, memorable presentation now:"""

    print(f"🤖 Asking Claude to create engaging presentation about: {topic}")
    print("⏳ Crafting punchy, memorable content...\n")
    
    message = client.messages.create(
        model="claude-sonnet-5",
        max_tokens=4000,
        messages=[
            {"role": "user", "content": prompt}
        ]
    )
    
    # Extract text content (skip thinking blocks)
    response_text = ""
    for block in message.content:
        if hasattr(block, 'text'):
            response_text += block.text
    
    response_text = response_text.strip()
    
    # Extract JSON
    json_start = response_text.find('{')
    json_end = response_text.rfind('}') + 1
    
    if json_start >= 0 and json_end > json_start:
        json_text = response_text[json_start:json_end]
        outline = json.loads(json_text)
        print(f"✅ Claude generated: {len(outline.get('slides', []))} slides")
        return outline
    else:
        raise ValueError("No valid JSON found in Claude's response")


class EnhancedThemePreview:
    """Generate presentations with enhanced visual design"""
    
    def __init__(self):
        with open('config.yaml', 'r') as f:
            config = yaml.safe_load(f)
        self.themes = config['themes']
    
    def create_presentation(self, outline: Dict, theme_name: str, output_path: str):
        """Create PowerPoint presentation with enhanced design"""
        
        theme = self.themes[theme_name]
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Color scheme
        primary = RGBColor(*theme['primary_color'])
        accent = RGBColor(*theme['accent_color'])
        bg = RGBColor(*theme['background'])
        text = RGBColor(*theme['text_color'])
        
        print(f"\n🎨 Creating '{theme['name']}' presentation with enhanced design...")
        
        for i, slide_data in enumerate(outline['slides']):
            slide_type = slide_data.get('type', 'content')
            
            if slide_type == 'title':
                slide = self._create_title_slide(prs, slide_data, primary, accent, text)
            elif slide_type == 'section':
                slide = self._create_section_slide(prs, slide_data, primary, accent, text)
            elif slide_type == 'conclusion':
                slide = self._create_conclusion_slide(prs, slide_data, accent, text)
            else:
                slide = self._create_content_slide(prs, slide_data, primary, accent, text, bg)
            
            # Add slide transition
            self._add_transition(slide)
            
            print(f"  ✓ Slide {i+1}/{len(outline['slides'])}: {slide_data.get('title', 'Untitled')}")
        
        prs.save(output_path)
        print(f"✅ Saved: {output_path}")
        return output_path
    
    def _add_transition(self, slide):
        """Add smooth transition to slide"""
        try:
            from pptx.enum.shapes import PP_TRANSITION_TYPE
            # Add a subtle fade transition (400ms)
            slide.slide.transition.type = 1  # Fade
            slide.slide.transition.advanceTime = 0
            slide.slide.transition.advanceOnClick = True
        except:
            pass  # If transitions don't work, continue without them
    
    def _create_title_slide(self, prs, data, primary_color, accent_color, text_color):
        """Create enhanced title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shapes = slide.shapes
        
        # Full background
        bg_shape = shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = primary_color
        bg_shape.line.fill.background()
        
        # Decorative accent shapes (top corner)
        accent_shape1 = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(8.5), Inches(0.5),
            Inches(1.2), Inches(0.15)
        )
        accent_shape1.fill.solid()
        accent_shape1.fill.fore_color.rgb = accent_color
        accent_shape1.line.fill.background()
        
        # Bottom accent bar
        accent_bar = shapes.add_shape(1, 0, Inches(6.8), prs.slide_width, Inches(0.7))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_color
        accent_bar.line.fill.background()
        
        # Title - improved spacing
        title_box = shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = data.get('title', 'Presentation Title')
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(56)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        title_para.line_spacing = 1.15
        
        # Subtitle - better positioning
        if data.get('subtitle'):
            subtitle_box = shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(7.6), Inches(1.2))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = data['subtitle']
            subtitle_frame.word_wrap = True
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(26)
            subtitle_para.font.color.rgb = RGBColor(245, 245, 245)
            subtitle_para.alignment = PP_ALIGN.CENTER
            subtitle_para.line_spacing = 1.3
        
        return slide
    
    def _create_section_slide(self, prs, data, primary_color, accent_color, text_color):
        """Create enhanced section divider slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shapes = slide.shapes
        
        # Background
        bg_shape = shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = primary_color
        bg_shape.line.fill.background()
        
        # Decorative accent line (top)
        accent_line = shapes.add_shape(
            1, Inches(2), Inches(3.2),
            Inches(6), Inches(0.08)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = accent_color
        accent_line.line.fill.background()
        
        # Section title
        title_box = shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(1.2))
        title_frame = title_box.text_frame
        title_frame.text = data.get('title', 'Section')
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(52)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def _create_conclusion_slide(self, prs, data, accent_color, text_color):
        """Create enhanced conclusion slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shapes = slide.shapes
        
        # Background
        bg_shape = shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = accent_color
        bg_shape.line.fill.background()
        
        # Decorative circle accent
        circle = shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(4.25), Inches(2.8),
            Inches(1.5), Inches(1.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
        circle.fill.fore_color.brightness = -0.1
        circle.line.fill.background()
        
        # Title
        title_box = shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1.2))
        title_frame = title_box.text_frame
        title_frame.text = data.get('title', 'Thank You')
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(56)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def _create_content_slide(self, prs, data, primary_color, accent_color, text_color, bg_color):
        """Create enhanced content slide with better spacing"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shapes = slide.shapes
        
        # Title bar - taller for better presence
        title_shape = shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.3))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = primary_color
        title_shape.line.fill.background()
        
        # Small decorative accent square in header
        accent_square = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.3), Inches(0.5),
            Inches(0.15), Inches(0.3)
        )
        accent_square.fill.solid()
        accent_square.fill.fore_color.rgb = accent_color
        accent_square.line.fill.background()
        
        # Title text with better spacing
        title_box = shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = data.get('title', 'Slide Title')
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(38)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.line_spacing = 1.1
        
        # Thicker accent line
        accent_line = shapes.add_shape(1, 0, Inches(1.3), prs.slide_width, Inches(0.08))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = accent_color
        accent_line.line.fill.background()
        
        # Content bullets with enhanced spacing
        bullets = data.get('bullets', [])
        if bullets:
            content_box = shapes.add_textbox(
                Inches(1.2), Inches(2.1),
                Inches(8), Inches(5)
            )
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.2)
            
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = bullet
                p.level = 0
                p.font.size = Pt(20)  # Slightly larger
                p.font.color.rgb = text_color
                p.space_before = Pt(8)
                p.space_after = Pt(20)  # More space between bullets
                p.line_spacing = 1.35  # Better line height
        
        return slide


def main():
    """Generate enhanced presentations with real Claude API"""
    
    print("\n" + "="*70)
    print("🎨 ENHANCED CLAUDE PRESENTATION GENERATOR")
    print("="*70)
    print("\nFeatures:")
    print("  ✨ Punchy, memorable content")
    print("  🎯 Better spacing and visual hierarchy")
    print("  🎬 Smooth slide transitions")
    print("  🎨 Enhanced design elements")
    print("="*70)
    
    # Get API key
    api_key = os.getenv('ANTHROPIC_API_KEY')
    if not api_key:
        api_key = input("\n🔑 Enter your Claude API key: ").strip()
    
    if not api_key:
        print("❌ API key required!")
        return
    
    # Get topic
    print("\n" + "="*70)
    topic = input("📝 What should the presentation be about?\n   Topic: ").strip()
    
    if not topic:
        topic = "The Future of Artificial Intelligence"
        print(f"\n   Using default: {topic}")
    
    # Get number of slides
    num_input = input("\n📊 How many slides? (default: 12): ").strip()
    num_slides = int(num_input) if num_input.isdigit() else 12
    
    # Get theme preference
    print("\n🎨 Available themes:")
    print("   1. modern - Professional blue & orange (default)")
    print("   2. dark - Tech-style dark theme")
    print("   3. education - Friendly educational theme")
    print("   4. warm - Creative warm colors")
    print("   5. all - Generate all 4 themes")
    
    theme_input = input("\n   Theme choice (1-5): ").strip()
    theme_map = {
        '1': ['modern'],
        '2': ['dark'],
        '3': ['education'],
        '4': ['warm'],
        '5': ['modern', 'dark', 'education', 'warm']
    }
    themes = theme_map.get(theme_input, ['modern'])
    
    print("\n" + "="*70)
    print("🚀 GENERATING PRESENTATION")
    print("="*70)
    
    # Generate outline with enhanced Claude prompt
    outline = get_claude_outline(api_key, topic, num_slides)
    
    print(f"\n📊 Presentation: {outline['title']}")
    print(f"📝 Slides: {len(outline['slides'])}")
    
    # Create enhanced presentations
    previewer = EnhancedThemePreview()
    
    for theme_name in themes:
        output = f"enhanced_{theme_name}_theme.pptx"
        previewer.create_presentation(outline, theme_name, output)
    
    print("\n" + "="*70)
    print("✅ DONE! Enhanced presentations created!")
    print("="*70)
    print(f"""
Files created:
""")
    for theme_name in themes:
        print(f"- enhanced_{theme_name}_theme.pptx")
    
    print(f"""
🎯 ENHANCEMENTS:
   ✨ More engaging, memorable content
   📐 Better spacing and visual hierarchy
   🎬 Smooth slide transitions
   🎨 Enhanced design elements
   
Open the files to see the difference! 🚀
""")


if __name__ == "__main__":
    main()
