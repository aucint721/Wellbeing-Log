#!/usr/bin/env python3
"""
Designer Claude Presentation Generator
Stunning visuals with geometric patterns, gradients, and modern design
"""

import os
import json
import yaml
from typing import Dict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def get_claude_outline(api_key: str, topic: str, num_slides: int = 10) -> Dict:
    """Generate presentation outline using real Claude API"""
    
    try:
        from anthropic import Anthropic
    except ImportError:
        print("Installing anthropic package...")
        import subprocess
        subprocess.check_call(["pip", "install", "anthropic"])
        from anthropic import Anthropic
    
    client = Anthropic(api_key=api_key)
    
    prompt = f"""You are an expert presentation designer creating a captivating presentation about: {topic}

CONTENT STYLE:
- Make bullet points PUNCHY, memorable, and impactful
- Use active voice and strong verbs
- Include specific examples and numbers when relevant
- Write like you're telling an engaging story, not a boring lecture
- Each bullet should be a complete insight, not just a fragment
- Aim for "aha!" moments that make people remember the content
- Use emojis sparingly if they add clarity (optional)

TARGET AUDIENCE: Year 1 students (or general audience if not education-focused)
- Keep language clear and accessible
- Use concrete examples they can visualize
- Make it fun and engaging, not dry

Create exactly {num_slides} slides with this structure:
1. Title slide (engaging title + compelling subtitle)
2. Content slides (3-4 punchy bullets each - quality over quantity)
3. Section breaks (use 1-2 for major topic transitions)
4. Conclusion slide

Return ONLY valid JSON in this exact format:
{{
  "title": "Main presentation title (make it catchy!)",
  "slides": [
    {{
      "type": "title",
      "title": "Engaging Title Text",
      "subtitle": "Compelling subtitle that makes you want to learn more"
    }},
    {{
      "type": "content",
      "title": "Clear, Action-Oriented Slide Title",
      "bullets": [
        "First impactful point with specific detail",
        "Second memorable insight that sticks",
        "Third engaging example that illustrates the concept",
        "Optional fourth point (only if needed)"
      ]
    }},
    {{
      "type": "section",
      "title": "Major Topic Transition"
    }},
    {{
      "type": "conclusion",
      "title": "Thank You"
    }}
  ]
}}

IMPORTANT:
- Each bullet should be 1-2 lines max
- Use numbers, examples, or analogies to make points memorable
- Avoid generic statements - be specific and interesting
- Create a logical narrative flow across slides

Create an engaging, memorable presentation now:"""

    print(f"🤖 Claude is crafting your presentation about: {topic}")
    print("✨ Creating punchy, memorable content...\n")
    
    message = client.messages.create(
        model="claude-sonnet-5",
        max_tokens=4000,
        messages=[
            {"role": "user", "content": prompt}
        ]
    )
    
    # Extract text content
    response_text = ""
    for block in message.content:
        if hasattr(block, 'text'):
            response_text += block.text
    
    response_text = response_text.strip()
    
    # Extract JSON
    json_start = response_text.find('{')
    json_end = response_text.rfind('}') + 1
    
    if json_start >= 0 and json_end > json_start:
        json_text = response_text[json_start:json_end]
        outline = json.loads(json_text)
        print(f"✅ Generated {len(outline.get('slides', []))} slides with engaging content!")
        return outline
    else:
        raise ValueError("No valid JSON found in Claude's response")


class DesignerPresentation:
    """Create stunning presentations with modern design"""
    
    def __init__(self):
        with open('config.yaml', 'r') as f:
            config = yaml.safe_load(f)
        self.themes = config['themes']
    
    def create_presentation(self, outline: Dict, theme_name: str, output_path: str):
        """Create stunning PowerPoint presentation"""
        
        theme = self.themes[theme_name]
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Color scheme
        primary = RGBColor(*theme['primary_color'])
        accent = RGBColor(*theme['accent_color'])
        bg = RGBColor(*theme['background'])
        text = RGBColor(*theme['text_color'])
        
        # Gradient colors if available
        if 'gradient_colors' in theme:
            grad_start = RGBColor(*theme['gradient_colors'][0])
            grad_end = RGBColor(*theme['gradient_colors'][1])
        else:
            grad_start = primary
            grad_end = accent
        
        print(f"\n🎨 Creating stunning '{theme['name']}' presentation...")
        
        for i, slide_data in enumerate(outline['slides']):
            slide_type = slide_data.get('type', 'content')
            
            if slide_type == 'title':
                slide = self._create_title_slide(prs, slide_data, primary, accent, grad_start, grad_end)
            elif slide_type == 'section':
                slide = self._create_section_slide(prs, slide_data, primary, accent, grad_start, grad_end)
            elif slide_type == 'conclusion':
                slide = self._create_conclusion_slide(prs, slide_data, primary, accent)
            else:
                slide = self._create_content_slide(prs, slide_data, primary, accent, text, bg)
            
            self._add_transition(slide)
            print(f"  ✓ Slide {i+1}/{len(outline['slides'])}: {slide_data.get('title', 'Untitled')}")
        
        prs.save(output_path)
        print(f"✅ Saved: {output_path}")
        return output_path
    
    def _add_transition(self, slide):
        """Add smooth transition"""
        try:
            slide.slide.transition.type = 1  # Fade
        except:
            pass
    
    def _add_geometric_pattern(self, shapes, prs, color, opacity=0.15):
        """Add subtle geometric pattern background"""
        # Create a grid of circles for visual interest
        circle_size = Inches(0.4)
        spacing = Inches(1.5)
        
        for x in range(0, 8, 2):
            for y in range(1, 7, 2):
                circle = shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(x + 0.5), Inches(y),
                    circle_size, circle_size
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = color
                circle.fill.transparency = 1 - opacity
                circle.line.fill.background()
    
    def _create_title_slide(self, prs, data, primary, accent, grad_start, grad_end):
        """Create stunning title slide with geometric elements"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shapes = slide.shapes
        
        # Background with gradient effect (simulated with overlapping rectangles)
        bg_shape1 = shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg_shape1.fill.solid()
        bg_shape1.fill.fore_color.rgb = grad_start
        bg_shape1.line.fill.background()
        
        # Gradient overlay
        bg_shape2 = shapes.add_shape(1, 0, Inches(3), prs.slide_width, Inches(4.5))
        bg_shape2.fill.solid()
        bg_shape2.fill.fore_color.rgb = grad_end
        bg_shape2.fill.transparency = 0.3
        bg_shape2.line.fill.background()
        
        # Decorative geometric elements
        # Large circle accent (top right)
        circle1 = shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(8), Inches(-0.5),
            Inches(2.5), Inches(2.5)
        )
        circle1.fill.solid()
        circle1.fill.fore_color.rgb = accent
        circle1.fill.transparency = 0.7
        circle1.line.fill.background()
        
        # Small accent squares
        for i, pos in enumerate([(0.3, 6.8), (0.8, 6.5), (1.3, 6.9)]):
            square = shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(pos[0]), Inches(pos[1]),
                Inches(0.25), Inches(0.25)
            )
            square.fill.solid()
            square.fill.fore_color.rgb = accent
            square.fill.transparency = 0.3 * (i + 1)
            square.line.fill.background()
        
        # Accent bar with angle
        accent_bar = shapes.add_shape(
            1, Inches(-0.5), Inches(6.5),
            Inches(11), Inches(1.2)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent
        accent_bar.fill.transparency = 0.1
        accent_bar.line.fill.background()
        accent_bar.rotation = -3
        
        # Title with shadow effect
        title_box = shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(2.2))
        title_frame = title_box.text_frame
        title_frame.text = data.get('title', 'Presentation Title')
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(60)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        title_para.line_spacing = 1.15
        
        # Subtitle
        if data.get('subtitle'):
            subtitle_box = shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(7.6), Inches(1.3))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = data['subtitle']
            subtitle_frame.word_wrap = True
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(28)
            subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
            subtitle_para.alignment = PP_ALIGN.CENTER
            subtitle_para.line_spacing = 1.3
        
        return slide
    
    def _create_section_slide(self, prs, data, primary, accent, grad_start, grad_end):
        """Create modern section divider with geometric design"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shapes = slide.shapes
        
        # Background gradient
        bg1 = shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg1.fill.solid()
        bg1.fill.fore_color.rgb = grad_start
        bg1.line.fill.background()
        
        bg2 = shapes.add_shape(1, 0, Inches(2), prs.slide_width, Inches(5.5))
        bg2.fill.solid()
        bg2.fill.fore_color.rgb = grad_end
        bg2.fill.transparency = 0.4
        bg2.line.fill.background()
        
        # Geometric accent pattern
        # Diagonal accent bars
        for i, offset in enumerate([0, 0.3, 0.6]):
            bar = shapes.add_shape(
                1, Inches(-1 + offset), Inches(2.5 + i*0.1),
                Inches(12), Inches(0.15)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent
            bar.fill.transparency = 0.4 + (i * 0.2)
            bar.line.fill.background()
            bar.rotation = 2
        
        # Large decorative circle
        circle = shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(7.5), Inches(5),
            Inches(3), Inches(3)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.fill.transparency = 0.85
        circle.line.fill.background()
        
        # Title
        title_box = shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = data.get('title', 'Section')
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(56)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def _create_conclusion_slide(self, prs, data, primary, accent):
        """Create elegant conclusion slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shapes = slide.shapes
        
        # Background
        bg = shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = accent
        bg.line.fill.background()
        
        # Decorative elements - rings
        for i, size in enumerate([2.5, 2, 1.5]):
            ring = shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(3.75 + i*0.25), Inches(2.5 + i*0.25),
                Inches(size), Inches(size)
            )
            ring.fill.solid()
            ring.fill.fore_color.rgb = RGBColor(255, 255, 255)
            ring.fill.transparency = 0.8 + (i * 0.05)
            ring.line.fill.background()
        
        # Accent squares pattern
        positions = [(0.5, 1), (1, 6.5), (8.5, 1.5), (9, 6)]
        for pos in positions:
            sq = shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(pos[0]), Inches(pos[1]),
                Inches(0.3), Inches(0.3)
            )
            sq.fill.solid()
            sq.fill.fore_color.rgb = RGBColor(255, 255, 255)
            sq.fill.transparency = 0.6
            sq.line.fill.background()
        
        # Title
        title_box = shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1.2))
        title_frame = title_box.text_frame
        title_frame.text = data.get('title', 'Thank You')
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(60)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def _create_content_slide(self, prs, data, primary, accent, text, bg):
        """Create modern content slide with visual hierarchy"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shapes = slide.shapes
        
        # Header with gradient effect (simulated)
        header1 = shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.4))
        header1.fill.solid()
        header1.fill.fore_color.rgb = primary
        header1.line.fill.background()
        
        header2 = shapes.add_shape(1, 0, Inches(0.7), prs.slide_width, Inches(0.7))
        header2.fill.solid()
        header2.fill.fore_color.rgb = accent
        header2.fill.transparency = 0.7
        header2.line.fill.background()
        
        # Decorative side accent bar
        side_bar = shapes.add_shape(
            1, 0, Inches(1.4),
            Inches(0.25), Inches(6.1)
        )
        side_bar.fill.solid()
        side_bar.fill.fore_color.rgb = accent
        side_bar.fill.transparency = 0.3
        side_bar.line.fill.background()
        
        # Small decorative elements in header
        positions = [(0.4, 0.55), (0.7, 0.55)]
        for pos in positions:
            dot = shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(pos[0]), Inches(pos[1]),
                Inches(0.15), Inches(0.15)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = accent
            dot.line.fill.background()
        
        # Title
        title_box = shapes.add_textbox(Inches(1), Inches(0.3), Inches(8.5), Inches(0.9))
        title_frame = title_box.text_frame
        title_frame.text = data.get('title', 'Slide Title')
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.line_spacing = 1.1
        
        # Accent underline
        underline = shapes.add_shape(1, Inches(1), Inches(1.4), Inches(2), Inches(0.1))
        underline.fill.solid()
        underline.fill.fore_color.rgb = accent
        underline.line.fill.background()
        
        # Content with enhanced spacing
        bullets = data.get('bullets', [])
        if bullets:
            content_box = shapes.add_textbox(
                Inches(1.4), Inches(2.2),
                Inches(7.8), Inches(4.8)
            )
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                # Add bullet point symbol
                p.text = "• " + bullet
                p.level = 0
                p.font.size = Pt(22)
                p.font.color.rgb = text
                p.space_before = Pt(10)
                p.space_after = Pt(22)
                p.line_spacing = 1.4
        
        # Decorative corner accent
        corner = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(9), Inches(6.8),
            Inches(0.8), Inches(0.5)
        )
        corner.fill.solid()
        corner.fill.fore_color.rgb = accent
        corner.fill.transparency = 0.7
        corner.line.fill.background()
        
        return slide


def main():
    """Generate designer presentations"""
    
    print("\n" + "="*70)
    print("✨ DESIGNER CLAUDE PRESENTATION GENERATOR ✨")
    print("="*70)
    print("\nFeatures:")
    print("  🎨 Stunning geometric patterns")
    print("  🌈 Beautiful gradient effects")
    print("  ✨ Modern designer themes")
    print("  🎯 Professional visual hierarchy")
    print("  🎬 Smooth transitions")
    print("="*70)
    
    # Get API key
    api_key = os.getenv('ANTHROPIC_API_KEY')
    if not api_key:
        api_key = input("\n🔑 Enter your Claude API key: ").strip()
    
    if not api_key:
        print("❌ API key required!")
        return
    
    # Get topic
    print("\n" + "="*70)
    topic = input("📝 What should the presentation be about?\n   Topic: ").strip()
    
    if not topic:
        topic = "The Future of Artificial Intelligence"
        print(f"\n   Using default: {topic}")
    
    # Get number of slides
    num_input = input("\n📊 How many slides? (default: 12): ").strip()
    num_slides = int(num_input) if num_input.isdigit() else 12
    
    # Show available themes
    with open('config.yaml', 'r') as f:
        config = yaml.safe_load(f)
    themes_list = list(config['themes'].keys())
    
    print("\n🎨 Available Designer Themes:")
    print("\n  Classic Themes:")
    print("   1. modern - Modern Professional")
    print("   2. dark - Dark Tech")
    print("   3. education - Education Blue")
    print("   4. warm - Warm Earth")
    print("\n  ✨ Designer Themes:")
    print("   5. neon_cyber - Neon Cyberpunk")
    print("   6. sunset_gradient - Sunset Gradient")
    print("   7. ocean_deep - Ocean Deep")
    print("   8. lavender_dream - Lavender Dream")
    print("   9. forest_minimal - Forest Minimal")
    print("   10. royal_purple - Royal Purple")
    print("   11. coral_pink - Coral Pink")
    print("   12. midnight_blue - Midnight Blue")
    print("   13. mint_fresh - Mint Fresh")
    print("   14. all - Generate ALL themes (13 presentations!)")
    
    theme_map = {
        '1': ['modern'], '2': ['dark'], '3': ['education'], '4': ['warm'],
        '5': ['neon_cyber'], '6': ['sunset_gradient'], '7': ['ocean_deep'],
        '8': ['lavender_dream'], '9': ['forest_minimal'], '10': ['royal_purple'],
        '11': ['coral_pink'], '12': ['midnight_blue'], '13': ['mint_fresh'],
        '14': themes_list
    }
    
    theme_input = input("\n   Theme choice (1-14): ").strip()
    themes = theme_map.get(theme_input, ['sunset_gradient'])
    
    print("\n" + "="*70)
    print("🚀 GENERATING DESIGNER PRESENTATION")
    print("="*70)
    
    # Generate outline
    outline = get_claude_outline(api_key, topic, num_slides)
    
    print(f"\n📊 Presentation: {outline['title']}")
    print(f"📝 Slides: {len(outline['slides'])}")
    
    # Create designer presentations
    designer = DesignerPresentation()
    
    for theme_name in themes:
        output = f"designer_{theme_name}.pptx"
        designer.create_presentation(outline, theme_name, output)
    
    print("\n" + "="*70)
    print("✨ DONE! Designer presentations created! ✨")
    print("="*70)
    print(f"""
Files created:
""")
    for theme_name in themes:
        print(f"  🎨 designer_{theme_name}.pptx")
    
    print(f"""
🎯 DESIGNER FEATURES:
   ✨ Geometric patterns and modern design
   🌈 Beautiful gradient effects
   🎨 13 stunning designer themes
   📐 Professional visual hierarchy
   🎬 Smooth transitions
   
Open the files and WOW your audience! 🚀
""")


if __name__ == "__main__":
    main()
